#!/usr/bin/env python3
"""
PowerPoint Generator - Convert markdown training content to branded .pptx
Applies Sparkry brand guidelines with minimal text on slides philosophy.
"""

import argparse
import re
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# Sparkry Brand Colors
SPARKRY_ORANGE = RGBColor(255, 107, 53)  # #ff6b35
SPARKRY_NAVY = RGBColor(23, 29, 40)      # #171d28
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)
DARK_GRAY = RGBColor(100, 100, 100)


class SlideContent:
    """Represents content for a single slide"""
    def __init__(self, slide_num: int, title: str, content: List[str], speaker_notes: str, slide_type: str = "content"):
        self.slide_num = slide_num
        self.title = title
        self.content = content
        self.speaker_notes = speaker_notes
        self.slide_type = slide_type  # "title", "section", "content", "stats"


class MarkdownParser:
    """Parse markdown training files into slide structure"""

    def __init__(self, markdown_path: Path):
        self.markdown_path = markdown_path
        self.content = markdown_path.read_text()

    def parse(self) -> Tuple[str, str, List[SlideContent]]:
        """Parse markdown and return (session_title, subtitle, slides)"""
        lines = self.content.split('\n')

        session_title = ""
        subtitle = ""
        slides = []

        current_slide = None
        current_content = []
        current_speaker_notes = []
        in_speaker_notes = False

        for line in lines:
            # Extract session title (# Session N: Title)
            if line.startswith('# ') and not session_title:
                session_title = line[2:].strip()
                continue

            # Extract subtitle (first paragraph after title)
            if not subtitle and session_title and line.strip() and not line.startswith('#'):
                subtitle = line.strip()
                continue

            # Section headers (## Section: Name)
            if line.startswith('## ') and 'Section:' in line:
                if current_slide:
                    current_slide.content = current_content
                    current_slide.speaker_notes = '\n'.join(current_speaker_notes)
                    slides.append(current_slide)

                section_name = line.split('Section:')[1].strip()
                current_slide = SlideContent(
                    slide_num=len(slides) + 1,
                    title=section_name,
                    content=[],
                    speaker_notes="",
                    slide_type="section"
                )
                current_content = []
                current_speaker_notes = []
                in_speaker_notes = False
                continue

            # Slide markers (### SLIDE N: Title)
            if line.startswith('### SLIDE'):
                if current_slide:
                    current_slide.content = current_content
                    current_slide.speaker_notes = '\n'.join(current_speaker_notes)
                    slides.append(current_slide)

                # Extract slide title
                match = re.match(r'### SLIDE (\d+): (.+)', line)
                if match:
                    slide_num = int(match.group(1))
                    title = match.group(2).strip()

                    # Determine slide type based on content indicators
                    slide_type = "content"
                    if any(keyword in title.lower() for keyword in ['reality', 'state', 'adoption']):
                        slide_type = "stats"

                    current_slide = SlideContent(
                        slide_num=slide_num,
                        title=title,
                        content=[],
                        speaker_notes="",
                        slide_type=slide_type
                    )
                    current_content = []
                    current_speaker_notes = []
                    in_speaker_notes = False
                continue

            # Speaker notes marker
            if '**SPEAKER NOTES:**' in line or 'SPEAKER NOTES:' in line:
                in_speaker_notes = True
                continue

            # Collect speaker notes
            if in_speaker_notes and current_slide:
                if line.strip() and not line.startswith('###'):
                    # Remove markdown formatting from speaker notes
                    clean_line = line.strip()
                    clean_line = re.sub(r'\*\*(.+?)\*\*', r'\1', clean_line)  # Remove bold
                    clean_line = re.sub(r'"(.+?)"', r'"\1"', clean_line)  # Preserve quotes
                    current_speaker_notes.append(clean_line)
                elif line.startswith('###'):
                    # Next slide starting, stop collecting notes
                    in_speaker_notes = False
                continue

            # Collect slide content (bullets)
            if current_slide and line.strip().startswith('-'):
                # Extract bullet text, removing markdown
                bullet = line.strip()[1:].strip()
                bullet = re.sub(r'\*\*(.+?)\*\*', r'\1', bullet)  # Remove bold
                bullet = re.sub(r'`(.+?)`', r'\1', bullet)  # Remove code formatting
                current_content.append(bullet)

        # Add last slide
        if current_slide:
            current_slide.content = current_content
            current_slide.speaker_notes = '\n'.join(current_speaker_notes)
            slides.append(current_slide)

        return session_title, subtitle, slides


class PresentationGenerator:
    """Generate PowerPoint presentation with Sparkry branding"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title: str, subtitle: str):
        """Create title slide with Sparkry branding"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Navy background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SPARKRY_NAVY

        # Title (large, white, with orange accent)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle (medium, light gray)
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = LIGHT_GRAY
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Orange accent bar at bottom
        accent = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(7.25), Inches(10), Inches(0.25)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = SPARKRY_ORANGE
        accent.line.fill.background()

    def add_section_slide(self, title: str):
        """Create section header slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Navy background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SPARKRY_NAVY

        # Section title (large, centered)
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = SPARKRY_ORANGE
        title_para.alignment = PP_ALIGN.CENTER

        # Orange accent line above
        accent = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(3), Inches(2.5), Inches(4), Inches(0.1)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = SPARKRY_ORANGE
        accent.line.fill.background()

    def add_content_slide(self, title: str, content: List[str], speaker_notes: str):
        """Create minimal text content slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # Title with orange accent
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = SPARKRY_NAVY

        # Orange underline under title
        accent = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1.4), Inches(2), Inches(0.08)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = SPARKRY_ORANGE
        accent.line.fill.background()

        # Content bullets (minimal - max 3, large text)
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(4)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # Only use first 3 bullets
            for bullet_text in content[:3]:
                para = content_frame.add_paragraph()
                para.text = bullet_text
                para.font.size = Pt(28)
                para.font.color.rgb = SPARKRY_NAVY
                para.level = 0
                para.space_before = Pt(12)
                para.space_after = Pt(12)

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def add_stats_slide(self, title: str, content: List[str], speaker_notes: str):
        """Create statistics slide with large numbers"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Navy background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SPARKRY_NAVY

        # Title (white, smaller)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE

        # Display statistics as large numbers
        y_position = 2.0
        for stat_text in content[:3]:  # Max 3 stats
            # Extract percentage/number if present
            match = re.search(r'(\d+%|\d+\.\d+%|\d+)', stat_text)
            if match:
                number = match.group(1)
                description = stat_text.replace(number, '').strip()

                # Large number in orange
                num_box = slide.shapes.add_textbox(
                    Inches(1), Inches(y_position), Inches(3), Inches(1)
                )
                num_frame = num_box.text_frame
                num_frame.text = number
                num_para = num_frame.paragraphs[0]
                num_para.font.size = Pt(72)
                num_para.font.bold = True
                num_para.font.color.rgb = SPARKRY_ORANGE

                # Description in white
                desc_box = slide.shapes.add_textbox(
                    Inches(4.5), Inches(y_position + 0.2), Inches(5), Inches(0.8)
                )
                desc_frame = desc_box.text_frame
                desc_frame.text = description
                desc_frame.word_wrap = True
                desc_para = desc_frame.paragraphs[0]
                desc_para.font.size = Pt(24)
                desc_para.font.color.rgb = WHITE

                y_position += 1.5
            else:
                # No number, just display text
                text_box = slide.shapes.add_textbox(
                    Inches(1), Inches(y_position), Inches(8), Inches(0.8)
                )
                text_frame = text_box.text_frame
                text_frame.text = stat_text
                text_frame.word_wrap = True
                text_para = text_frame.paragraphs[0]
                text_para.font.size = Pt(28)
                text_para.font.color.rgb = WHITE

                y_position += 1.2

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def generate(self, session_title: str, subtitle: str, slides: List[SlideContent], output_path: Path):
        """Generate complete presentation"""
        # Add title slide
        self.add_title_slide(session_title, subtitle)

        # Add content slides
        for slide in slides:
            if slide.slide_type == "section":
                self.add_section_slide(slide.title)
            elif slide.slide_type == "stats":
                self.add_stats_slide(slide.title, slide.content, slide.speaker_notes)
            else:
                self.add_content_slide(slide.title, slide.content, slide.speaker_notes)

        # Save presentation
        self.prs.save(str(output_path))
        print(f"✅ Generated presentation: {output_path}")
        print(f"   - {len(slides) + 1} slides total")
        print(f"   - Speaker notes embedded")


def main():
    parser = argparse.ArgumentParser(
        description="Generate branded PowerPoint presentation from markdown"
    )
    parser.add_argument(
        "--input",
        required=True,
        type=Path,
        help="Input markdown file path"
    )
    parser.add_argument(
        "--output",
        required=True,
        type=Path,
        help="Output .pptx file path"
    )

    args = parser.parse_args()

    # Validate input file exists
    if not args.input.exists():
        print(f"❌ Error: Input file not found: {args.input}")
        return 1

    # Create output directory if needed
    args.output.parent.mkdir(parents=True, exist_ok=True)

    # Parse markdown
    print(f"📖 Parsing markdown: {args.input}")
    parser = MarkdownParser(args.input)
    session_title, subtitle, slides = parser.parse()

    print(f"   - Session: {session_title}")
    print(f"   - Found {len(slides)} slides")

    # Generate presentation
    print("🎨 Generating presentation with Sparkry branding...")
    generator = PresentationGenerator()
    generator.generate(session_title, subtitle, slides, args.output)

    return 0


if __name__ == "__main__":
    exit(main())
