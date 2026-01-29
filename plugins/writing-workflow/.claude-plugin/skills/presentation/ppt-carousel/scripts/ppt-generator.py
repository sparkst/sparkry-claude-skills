#!/usr/bin/env python3
"""
PowerPoint Generator

Generate PowerPoint carousel presentation from slide manifest with brand styling.

Usage:
    python ppt-generator.py slides.json --background bg.png --format square --output carousel.pptx

Output:
    Creates .pptx file with brand-compliant slides
"""

import json
import sys
from pathlib import Path
from typing import Dict, Any, List, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print(json.dumps({
        'error': 'python-pptx library not installed. Run: pip install python-pptx'
    }), file=sys.stderr)
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print(json.dumps({
        'error': 'Pillow library not installed. Run: pip install Pillow'
    }), file=sys.stderr)
    sys.exit(1)


# Sparkry Brand Colors (RGB)
COLORS = {
    'primary_orange': RGBColor(255, 107, 53),    # #ff6b35
    'navy': RGBColor(23, 29, 40),                 # #171d28
    'electric_blue': RGBColor(14, 165, 233),      # #0ea5e9
    'electric_cyan': RGBColor(0, 217, 255),       # #00d9ff
    'white': RGBColor(255, 255, 255),             # #ffffff
    'text_primary': RGBColor(15, 23, 42),         # #0f172a
    'text_muted': RGBColor(100, 116, 139),        # #64748b
    'muted_bg': RGBColor(241, 245, 249),          # #f1f5f9
}

# Format dimensions (width, height in inches)
FORMATS = {
    'square': (10, 10),      # 1080×1080px at 108 DPI
    'portrait': (10, 12.5),  # 1080×1350px at 108 DPI
}

# Layout margins (inches)
MARGINS = {
    'slide': 0.74,      # 80px at 108 DPI
    'text_block': 0.37,  # 40px at 108 DPI
}

# Logo protection zone (bottom-right, inches)
LOGO_ZONE = {
    'width': 1.85,   # 200px at 108 DPI
    'height': 1.20,  # 130px at 108 DPI
}


def create_presentation(format_name: str = 'square') -> Presentation:
    """
    Create new PowerPoint presentation with specified dimensions.

    Args:
        format_name: Format (square or portrait)

    Returns:
        Presentation object
    """
    prs = Presentation()

    # Set slide dimensions
    width, height = FORMATS.get(format_name, FORMATS['square'])
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    return prs


def apply_background(slide, background_path: Optional[str] = None):
    """
    Apply background image to slide.

    Args:
        slide: Slide object
        background_path: Path to background image (optional)
    """
    if not background_path or not Path(background_path).exists():
        # Use solid navy background as default
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['navy']
        return

    # Add background image
    try:
        left = Inches(0)
        top = Inches(0)
        width = slide.width
        height = slide.height

        # Add image as background (behind all content)
        slide.shapes.add_picture(
            background_path,
            left, top,
            width=width,
            height=height
        )
    except Exception as e:
        print(f"Warning: Failed to apply background image: {e}", file=sys.stderr)


def add_text_box(slide, text: str, left: float, top: float, width: float, height: float,
                 font_name: str = 'Inter', font_size: int = 28, font_color: RGBColor = None,
                 bold: bool = False, alignment: PP_ALIGN = PP_ALIGN.LEFT) -> Any:
    """
    Add text box to slide with formatting.

    Args:
        slide: Slide object
        text: Text content
        left, top, width, height: Position and size (inches)
        font_name: Font name
        font_size: Font size (pt)
        font_color: RGB color
        bold: Bold weight
        alignment: Text alignment

    Returns:
        TextBox shape
    """
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )

    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    # Format paragraph
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignment

    # Format font
    font = paragraph.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold

    if font_color:
        font.color.rgb = font_color

    return textbox


def add_icon(slide, icon_path: str, left: float, top: float, size: float):
    """
    Add icon image to slide.

    Args:
        slide: Slide object
        icon_path: Path to icon image
        left, top: Position (inches)
        size: Icon size (inches)
    """
    if not Path(icon_path).exists():
        print(f"Warning: Icon not found: {icon_path}", file=sys.stderr)
        return

    try:
        slide.shapes.add_picture(
            icon_path,
            Inches(left), Inches(top),
            width=Inches(size),
            height=Inches(size)
        )
    except Exception as e:
        print(f"Warning: Failed to add icon: {e}", file=sys.stderr)


def create_hook_slide(prs: Presentation, slide_data: Dict[str, Any], background_path: Optional[str],
                     icon_cache_dir: Path) -> None:
    """
    Create hook slide (opening, attention-grabbing).

    Layout:
    - Center-aligned
    - Large icon at top
    - Title + subtitle
    - Gradient overlay for contrast
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Apply background
    apply_background(slide, background_path)

    # Add semi-transparent overlay for text contrast
    # (Note: python-pptx doesn't support gradients easily, use solid overlay)

    # Icon (top-center)
    if slide_data.get('icon'):
        icon_name = slide_data['icon'].replace(':', '-')
        icon_path = icon_cache_dir / f"{icon_name}-ff6b35-100.png"
        if icon_path.exists():
            icon_size = 0.93  # 100px at 108 DPI
            icon_left = (prs.slide_width / Inches(1) - icon_size) / 2
            add_icon(slide, str(icon_path), icon_left, 1.5, icon_size)

    # Title (center)
    title = slide_data.get('title', '')
    title_top = 3.5 if slide_data.get('icon') else 2.5
    add_text_box(
        slide, title,
        left=MARGINS['slide'], top=title_top,
        width=prs.slide_width / Inches(1) - 2 * MARGINS['slide'],
        height=2,
        font_name='Poppins', font_size=48, font_color=COLORS['white'],
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # Subtitle (center)
    subtitle = slide_data.get('content', '')
    if subtitle:
        add_text_box(
            slide, subtitle,
            left=MARGINS['slide'], top=title_top + 1.5,
            width=prs.slide_width / Inches(1) - 2 * MARGINS['slide'],
            height=2,
            font_name='Inter', font_size=36, font_color=COLORS['white'],
            bold=False, alignment=PP_ALIGN.CENTER
        )


def create_framework_slide(prs: Presentation, slide_data: Dict[str, Any], background_path: Optional[str],
                          icon_cache_dir: Path) -> None:
    """
    Create framework slide (core content with bullets).

    Layout:
    - Left-aligned
    - Icon top-left
    - Title + numbered bullets
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Apply background
    apply_background(slide, background_path)

    # Icon (top-left)
    if slide_data.get('icon'):
        icon_name = slide_data['icon'].replace(':', '-')
        icon_path = icon_cache_dir / f"{icon_name}-ff6b35-80.png"
        if icon_path.exists():
            add_icon(slide, str(icon_path), MARGINS['slide'], MARGINS['slide'], 0.74)  # 80px

    # Title (left-aligned, next to icon)
    title = slide_data.get('title', '')
    add_text_box(
        slide, title,
        left=MARGINS['slide'] + 1,  # After icon
        top=MARGINS['slide'],
        width=prs.slide_width / Inches(1) - MARGINS['slide'] - 1 - MARGINS['slide'],
        height=1,
        font_name='Poppins', font_size=48, font_color=COLORS['white'],
        bold=True, alignment=PP_ALIGN.LEFT
    )

    # Bullets (left-aligned)
    content = slide_data.get('content', [])
    if isinstance(content, str):
        content = [content]

    bullet_top = MARGINS['slide'] + 1.5
    bullet_spacing = 0.8

    for i, bullet in enumerate(content[:5]):  # Max 5 bullets
        # Bullet text
        bullet_text = bullet.strip('*→ ').strip()
        add_text_box(
            slide, bullet_text,
            left=MARGINS['slide'], top=bullet_top + i * bullet_spacing,
            width=prs.slide_width / Inches(1) - 2 * MARGINS['slide'] - LOGO_ZONE['width'],
            height=bullet_spacing,
            font_name='Inter', font_size=28, font_color=COLORS['white'],
            bold=False, alignment=PP_ALIGN.LEFT
        )


def create_example_slide(prs: Presentation, slide_data: Dict[str, Any], background_path: Optional[str],
                        icon_cache_dir: Path) -> None:
    """
    Create example slide (real-world story).

    Layout:
    - "REAL EXAMPLE" label
    - Icon top-right
    - Story text left-aligned
    - Light background tint
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Apply background (lighter tint for examples)
    if background_path:
        apply_background(slide, background_path)
    else:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['muted_bg']

    # "REAL EXAMPLE" label (top-left)
    add_text_box(
        slide, "REAL EXAMPLE",
        left=MARGINS['slide'], top=MARGINS['slide'],
        width=4, height=0.3,
        font_name='Poppins', font_size=12, font_color=COLORS['primary_orange'],
        bold=True, alignment=PP_ALIGN.LEFT
    )

    # Icon (top-right)
    if slide_data.get('icon'):
        icon_name = slide_data['icon'].replace(':', '-')
        icon_path = icon_cache_dir / f"{icon_name}-ff6b35-60.png"
        if icon_path.exists():
            icon_left = prs.slide_width / Inches(1) - MARGINS['slide'] - 0.56  # 60px
            add_icon(slide, str(icon_path), icon_left, MARGINS['slide'], 0.56)

    # Content text (left-aligned)
    content = slide_data.get('content', '')
    add_text_box(
        slide, content,
        left=MARGINS['slide'], top=MARGINS['slide'] + 1,
        width=prs.slide_width / Inches(1) - 2 * MARGINS['slide'],
        height=6,
        font_name='Inter', font_size=32, font_color=COLORS['text_primary'],
        bold=False, alignment=PP_ALIGN.LEFT
    )


def create_diagnostic_slide(prs: Presentation, slide_data: Dict[str, Any], background_path: Optional[str]) -> None:
    """
    Create diagnostic slide (key question/test).

    Layout:
    - Navy background
    - Center-aligned
    - Large text
    - High contrast
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Navy background (solid)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['navy']

    # Content text (center-aligned, large)
    content = slide_data.get('content', '')
    add_text_box(
        slide, content,
        left=MARGINS['slide'], top=3,
        width=prs.slide_width / Inches(1) - 2 * MARGINS['slide'],
        height=4,
        font_name='Poppins', font_size=40, font_color=COLORS['white'],
        bold=True, alignment=PP_ALIGN.CENTER
    )


def create_cta_slide(prs: Presentation, slide_data: Dict[str, Any], background_path: Optional[str],
                    icon_cache_dir: Path) -> None:
    """
    Create CTA slide (call-to-action).

    Layout:
    - Center-aligned
    - Icon center
    - Question + action
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Apply background
    apply_background(slide, background_path)

    # Icon (center)
    if slide_data.get('icon'):
        icon_name = slide_data['icon'].replace(':', '-')
        icon_path = icon_cache_dir / f"{icon_name}-ff6b35-80.png"
        if icon_path.exists():
            icon_size = 0.74
            icon_left = (prs.slide_width / Inches(1) - icon_size) / 2
            add_icon(slide, str(icon_path), icon_left, 2, icon_size)

    # Question (center)
    title = slide_data.get('title', '')
    add_text_box(
        slide, title,
        left=MARGINS['slide'], top=3.5,
        width=prs.slide_width / Inches(1) - 2 * MARGINS['slide'],
        height=2,
        font_name='Poppins', font_size=40, font_color=COLORS['white'],
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # Action (center, orange)
    content = slide_data.get('content', '')
    add_text_box(
        slide, content,
        left=MARGINS['slide'], top=5.5,
        width=prs.slide_width / Inches(1) - 2 * MARGINS['slide'],
        height=1.5,
        font_name='Inter', font_size=32, font_color=COLORS['primary_orange'],
        bold=False, alignment=PP_ALIGN.CENTER
    )


def generate_presentation(slides_manifest: Dict[str, Any], background_path: Optional[str],
                         format_name: str, output_path: str) -> Dict[str, Any]:
    """
    Generate PowerPoint presentation from slide manifest.

    Args:
        slides_manifest: Slide data from optimizer
        background_path: Path to background image
        format_name: Format (square or portrait)
        output_path: Output .pptx path

    Returns:
        Dict with generation summary
    """
    # Create presentation
    prs = create_presentation(format_name)

    slides = slides_manifest.get('slides', [])
    icon_cache_dir = Path(__file__).parent.parent / 'cache' / 'icons'

    print(f"✓ Presentation initialized: {FORMATS[format_name][0]}×{FORMATS[format_name][1]} inches ({format_name})", file=sys.stderr)

    if background_path:
        print(f"✓ Background: {background_path}", file=sys.stderr)

    # Generate slides
    for i, slide_data in enumerate(slides):
        slide_type = slide_data.get('type', 'framework')

        print(f"✓ Slide {i+1}/{len(slides)}: {slide_type.title()} (icon: {slide_data.get('icon', 'none')})", file=sys.stderr)

        if slide_type == 'hook':
            create_hook_slide(prs, slide_data, background_path, icon_cache_dir)
        elif slide_type == 'framework':
            create_framework_slide(prs, slide_data, background_path, icon_cache_dir)
        elif slide_type == 'example':
            create_example_slide(prs, slide_data, background_path, icon_cache_dir)
        elif slide_type == 'diagnostic':
            create_diagnostic_slide(prs, slide_data, background_path)
        elif slide_type == 'cta':
            create_cta_slide(prs, slide_data, background_path, icon_cache_dir)

    # Save presentation
    prs.save(output_path)

    # Get file size
    file_size_mb = Path(output_path).stat().st_size / (1024 * 1024)

    print(f"✓ File saved: {output_path} ({len(slides)} slides, {file_size_mb:.1f} MB)", file=sys.stderr)

    return {
        'success': True,
        'output': output_path,
        'slides': len(slides),
        'format': format_name,
        'file_size_mb': round(file_size_mb, 1)
    }


def main():
    if len(sys.argv) < 2:
        print("Usage: python ppt-generator.py <slides.json> [--background bg.png] [--format square] [--output carousel.pptx]", file=sys.stderr)
        sys.exit(1)

    slides_file = sys.argv[1]

    # Parse arguments
    background_path = None
    format_name = 'square'
    output_path = 'carousel.pptx'

    for i, arg in enumerate(sys.argv):
        if arg == '--background' and i + 1 < len(sys.argv):
            background_path = sys.argv[i + 1]
        elif arg == '--format' and i + 1 < len(sys.argv):
            format_name = sys.argv[i + 1]
        elif arg == '--output' and i + 1 < len(sys.argv):
            output_path = sys.argv[i + 1]

    # Validate files
    if not Path(slides_file).exists():
        print(json.dumps({
            'error': f"Slides file not found: {slides_file}"
        }))
        sys.exit(1)

    if background_path and not Path(background_path).exists():
        print(f"Warning: Background image not found: {background_path}", file=sys.stderr)
        background_path = None

    try:
        # Load slides manifest
        with open(slides_file, 'r', encoding='utf-8') as f:
            slides_manifest = json.load(f)

        # Generate presentation
        result = generate_presentation(slides_manifest, background_path, format_name, output_path)

        # Output JSON summary
        print(json.dumps(result, indent=2))

    except Exception as e:
        print(json.dumps({
            'error': str(e)
        }), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
