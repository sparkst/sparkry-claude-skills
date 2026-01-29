#!/usr/bin/env python3
"""
Brand Validator

Validate PowerPoint presentation for Sparkry brand compliance.
Checks colors, fonts, spacing, logo zone, and readability.

Usage:
    python brand-validator.py carousel.pptx

Output (JSON):
    {
      "valid": true,
      "checks": [...],
      "warnings": 1,
      "errors": 0
    }
"""

import json
import sys
from pathlib import Path
from typing import Dict, Any, List

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print(json.dumps({
        'error': 'python-pptx library not installed. Run: pip install python-pptx'
    }), file=sys.stderr)
    sys.exit(1)


# Sparkry Brand Colors (RGB tuples)
BRAND_COLORS = {
    (255, 107, 53): 'primary_orange',    # #ff6b35
    (23, 29, 40): 'navy',                # #171d28
    (14, 165, 233): 'electric_blue',     # #0ea5e9
    (0, 217, 255): 'electric_cyan',      # #00d9ff
    (255, 255, 255): 'white',            # #ffffff
    (15, 23, 42): 'text_primary',        # #0f172a
    (100, 116, 139): 'text_muted',       # #64748b
    (241, 245, 249): 'muted_bg',         # #f1f5f9
}

# Approved fonts
APPROVED_FONTS = ['Poppins', 'Inter', 'Arial', 'Helvetica', 'sans-serif']

# Logo protection zone (bottom-right, inches)
LOGO_ZONE = {
    'width': 1.85,   # 200px at 108 DPI
    'height': 1.20,  # 130px at 108 DPI
}

# Minimum margins (inches)
MIN_MARGIN = 0.74  # 80px at 108 DPI

# Readability standards
MAX_LINES_PER_SLIDE = 5
MIN_BODY_FONT_SIZE = 28  # pt


def rgb_color_approx(color_tuple: tuple, tolerance: int = 10) -> bool:
    """
    Check if color is approximately in brand palette.

    Args:
        color_tuple: RGB tuple (r, g, b)
        tolerance: Color matching tolerance

    Returns:
        True if color is within tolerance of brand colors
    """
    if color_tuple in BRAND_COLORS:
        return True

    # Check with tolerance
    for brand_color in BRAND_COLORS.keys():
        if all(abs(c1 - c2) <= tolerance for c1, c2 in zip(color_tuple, brand_color)):
            return True

    return False


def validate_colors(prs: Presentation) -> Dict[str, Any]:
    """
    Validate that all colors are within Sparkry brand palette.

    Args:
        prs: Presentation object

    Returns:
        Validation result dict
    """
    non_brand_colors = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check text colors
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.color and run.font.color.type == 1:  # RGB color
                            rgb = (
                                run.font.color.rgb[0],
                                run.font.color.rgb[1],
                                run.font.color.rgb[2]
                            )
                            if not rgb_color_approx(rgb):
                                non_brand_colors.append({
                                    'slide': slide_idx + 1,
                                    'color': f"rgb{rgb}",
                                    'location': 'text'
                                })

    if non_brand_colors:
        return {
            'check': 'color_palette',
            'status': 'warning',
            'message': f"Found {len(non_brand_colors)} colors outside brand palette",
            'details': non_brand_colors[:5]  # First 5 only
        }

    return {
        'check': 'color_palette',
        'status': 'pass',
        'message': 'All colors within Sparkry palette'
    }


def validate_fonts(prs: Presentation) -> Dict[str, Any]:
    """
    Validate that approved fonts are used.

    Args:
        prs: Presentation object

    Returns:
        Validation result dict
    """
    non_approved_fonts = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name and run.font.name not in APPROVED_FONTS:
                            non_approved_fonts.append({
                                'slide': slide_idx + 1,
                                'font': run.font.name
                            })

    if non_approved_fonts:
        # Check if fallback fonts are used (acceptable warning)
        fallback_only = all(f['font'] in ['Arial', 'Helvetica'] for f in non_approved_fonts)

        if fallback_only:
            return {
                'check': 'font_compliance',
                'status': 'warning',
                'message': 'Poppins/Inter not found, fallback to Arial/Helvetica'
            }

        return {
            'check': 'font_compliance',
            'status': 'warning',
            'message': f"Found {len(non_approved_fonts)} non-approved fonts",
            'details': non_approved_fonts[:5]
        }

    return {
        'check': 'font_compliance',
        'status': 'pass',
        'message': 'All fonts approved (Poppins/Inter)'
    }


def validate_logo_zone(prs: Presentation) -> Dict[str, Any]:
    """
    Validate that logo protection zone (bottom-right) is clear.

    Args:
        prs: Presentation object

    Returns:
        Validation result dict
    """
    violations = []

    slide_width = prs.slide_width / Inches(1)
    slide_height = prs.slide_height / Inches(1)

    # Logo zone boundaries (bottom-right)
    logo_left = slide_width - LOGO_ZONE['width']
    logo_top = slide_height - LOGO_ZONE['height']

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Get shape position and size (inches)
            shape_left = shape.left / Inches(1)
            shape_top = shape.top / Inches(1)
            shape_right = shape_left + (shape.width / Inches(1))
            shape_bottom = shape_top + (shape.height / Inches(1))

            # Check if shape overlaps logo zone
            if (shape_right > logo_left and shape_bottom > logo_top):
                violations.append({
                    'slide': slide_idx + 1,
                    'shape_type': shape.shape_type
                })

    if violations:
        return {
            'check': 'logo_zone',
            'status': 'error',
            'message': f"Logo zone violated on {len(violations)} slides",
            'details': violations
        }

    return {
        'check': 'logo_zone',
        'status': 'pass',
        'message': 'Bottom-right 200×130px clear on all slides'
    }


def validate_spacing(prs: Presentation) -> Dict[str, Any]:
    """
    Validate that proper margins are maintained.

    Args:
        prs: Presentation object

    Returns:
        Validation result dict
    """
    margin_violations = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check left/top margins
            shape_left = shape.left / Inches(1)
            shape_top = shape.top / Inches(1)

            if shape_left < MIN_MARGIN or shape_top < MIN_MARGIN:
                margin_violations.append({
                    'slide': slide_idx + 1,
                    'issue': 'insufficient_margin'
                })

    if margin_violations:
        return {
            'check': 'spacing',
            'status': 'warning',
            'message': f"Margin violations on {len(margin_violations)} slides",
            'details': margin_violations[:5]
        }

    return {
        'check': 'spacing',
        'status': 'pass',
        'message': 'Margins 80-100px maintained'
    }


def validate_readability(prs: Presentation) -> Dict[str, Any]:
    """
    Validate mobile readability standards.

    Args:
        prs: Presentation object

    Returns:
        Validation result dict
    """
    issues = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                # Count lines
                line_count = len(shape.text_frame.paragraphs)

                if line_count > MAX_LINES_PER_SLIDE:
                    issues.append({
                        'slide': slide_idx + 1,
                        'issue': f'too_many_lines ({line_count} > {MAX_LINES_PER_SLIDE})'
                    })

                # Check font sizes
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size < Pt(MIN_BODY_FONT_SIZE):
                            issues.append({
                                'slide': slide_idx + 1,
                                'issue': f'font_too_small ({run.font.size.pt}pt < {MIN_BODY_FONT_SIZE}pt)'
                            })

    if issues:
        return {
            'check': 'readability',
            'status': 'warning',
            'message': f"Readability issues on {len(issues)} slides",
            'details': issues[:5]
        }

    return {
        'check': 'readability',
        'status': 'pass',
        'message': 'Mobile readability optimized'
    }


def validate_presentation(pptx_path: str) -> Dict[str, Any]:
    """
    Run all brand validation checks on presentation.

    Args:
        pptx_path: Path to .pptx file

    Returns:
        Validation report dict
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {
            'valid': False,
            'error': f"Failed to load presentation: {e}"
        }

    # Run all checks
    checks = [
        validate_colors(prs),
        validate_fonts(prs),
        validate_logo_zone(prs),
        validate_spacing(prs),
        validate_readability(prs),
    ]

    # Count warnings and errors
    warnings = sum(1 for c in checks if c['status'] == 'warning')
    errors = sum(1 for c in checks if c['status'] == 'error')

    # Overall validity (no errors, warnings acceptable)
    valid = errors == 0

    return {
        'valid': valid,
        'checks': checks,
        'warnings': warnings,
        'errors': errors
    }


def main():
    if len(sys.argv) < 2:
        print("Usage: python brand-validator.py <carousel.pptx>", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]

    if not Path(pptx_path).exists():
        print(json.dumps({
            'error': f"File not found: {pptx_path}"
        }))
        sys.exit(1)

    try:
        # Validate presentation
        result = validate_presentation(pptx_path)

        # Output JSON
        print(json.dumps(result, indent=2))

        # Exit with error code if validation failed
        if not result.get('valid', False):
            sys.exit(1)

    except Exception as e:
        print(json.dumps({
            'error': str(e)
        }), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
